import os
import io
import requests
from bs4 import BeautifulSoup
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain_community.llms import Cohere
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import PromptTemplate
import tempfile
import urllib.request
from typing import List, Optional, Dict, Any, Union
import re

class RAGManager:
    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        embedding_model: str = 'sentence-transformers/all-MiniLM-L6-v2'
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.embedding_model = embedding_model
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            separators=['\n', '\n\n', ' ', '']
        )
        
        self.embeddings = HuggingFaceEmbeddings(model_name=self.embedding_model)
        
        self.documents = []
        self.extracted_text = ""
        self.chunks = []
        self.vectorstore = None
        self.retriever = None

    def extract_text_from_pdf(self, file_path: str) -> str:
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"Error extracting text from PDF {file_path}: {e}")
            return ""

    def extract_text_from_docx(self, file_path: str) -> str:
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            print(f"Error extracting text from DOCX {file_path}: {e}")
            return ""

    def extract_text_from_pptx(self, file_path: str) -> str:
        try:
            ppt = Presentation(file_path)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error extracting text from PPTX {file_path}: {e}")
            return ""

    def extract_text_from_url(self, url: str) -> str:
        try:
            response = requests.get(url)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            for script in soup(["script", "style"]):
                script.extract()
                
            text = soup.get_text(separator="\n")
            
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        except Exception as e:
            print(f"Error extracting text from URL {url}: {e}")
            return ""

    def load_document(self, source: str) -> None:
        text = ""
        
        if source.startswith(('http://', 'https://')):
            if source.endswith(('.pdf', '.docx', '.pptx')):
                file_path = self.download_file(source)
                if file_path:
                    self.documents.append(file_path)
                    if source.endswith('.pdf'):
                        text = self.extract_text_from_pdf(file_path)
                    elif source.endswith('.docx'):
                        text = self.extract_text_from_docx(file_path)
                    elif source.endswith('.pptx'):
                        text = self.extract_text_from_pptx(file_path)
            else:
                text = self.extract_text_from_url(source)
        else:
            self.documents.append(source)
            if source.endswith('.pdf'):
                text = self.extract_text_from_pdf(source)
            elif source.endswith('.docx'):
                text = self.extract_text_from_docx(source)
            elif source.endswith('.pptx'):
                text = self.extract_text_from_pptx(source)
            else:
                print(f"Unsupported file format: {source}")
                
        if text:
            self.extracted_text += text + "\n\n"
            print(f"Successfully loaded document: {source} ({len(text)} characters)")
        else:
            print(f"Failed to extract text from: {source}")

    def process_text(self, text: Optional[str] = None) -> None:
        if text is None:
            text = self.extracted_text
            
        if not text:
            print("No text to process. Please load documents first.")
            return
            
        self.chunks = self.text_splitter.split_text(text=text)
        print(f"Created {len(self.chunks)} text chunks")

    def build_vectorstore(self) -> None:
        if not self.chunks:
            print("No chunks to index. Please process text first.")
            return
            
        self.vectorstore = FAISS.from_texts(self.chunks, embedding=self.embeddings)
        self.retriever = self.vectorstore.as_retriever(
            search_type="similarity", 
            search_kwargs={"k": 6}
        )
        print(f"Built vector store with {len(self.chunks)} chunks")

    def save_vectorstore(self, path: str) -> None:
        if self.vectorstore:
            self.vectorstore.save_local(path)
            print(f"Vector store saved to {path}")
        else:
            print("No vector store to save. Please build vector store first.")

    def load_vectorstore(self, path: str) -> None:
        try:
            self.vectorstore = FAISS.load_local(
                path, 
                self.embeddings,
                allow_dangerous_deserialization=True
            )
            self.retriever = self.vectorstore.as_retriever(
                search_type="similarity", 
                search_kwargs={"k": 6}
            )
            print(f"Vector store loaded from {path}")
        except Exception as e:
            print(f"Error loading vector store: {e}")
            if self.chunks:
                print("Attempting to rebuild vector store from chunks...")
                self.build_vectorstore()
                self.save_vectorstore(path)
            else:
                print("No chunks available to rebuild vector store")

    def format_docs(self, docs):
        return "\n\n".join(doc.page_content for doc in docs)

    def query(self, question: str, model: str = "command", temperature: float = 0.1) -> str:
        if not self.retriever:
            print("No retriever available. Please build vector store first.")
            return "Error: Vector store not built"
            
        prompt_template = """Answer the question as precise as possible using the provided context. 
                          If the answer is not contained in the context, say "answer not available in context" \n\n
                          Context: \n {context}?\n
                          Question: \n {question} \n
                          Answer:"""
        
        prompt = PromptTemplate.from_template(template=prompt_template)
        
        cohere_llm = Cohere(
            model=model, 
            temperature=temperature, 
            cohere_api_key=os.getenv('COHERE_API_KEY')
        )
        
        def get_context(query):
            docs = self.retriever.invoke(query)
            return self.format_docs(docs)
            
        rag_chain = (
            {"context": RunnablePassthrough() | get_context, "question": RunnablePassthrough()}
            | prompt
            | cohere_llm
            | StrOutputParser()
        )
        
        return rag_chain.invoke(question) 